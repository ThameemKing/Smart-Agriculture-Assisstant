from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
BRAND_COLOR = RGBColor(34, 139, 34)  # Forest Green
ACCENT_COLOR = RGBColor(255, 140, 0)  # Dark Orange
TEXT_COLOR = RGBColor(40, 40, 40)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BRAND_COLOR
    title_shape.line.color.rgb = BRAND_COLOR
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    title_frame.margin_top = Inches(0.15)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "AI FARMING ASSISTANT", "Empowering Indian Farmers with Offline AI Technology")

# Slide 2: The Problem
add_content_slide(prs, "The Problem: Market Gap", [
    "🚜 80+ million small farmers in India lack access to agricultural advisory",
    "📱 Poor internet connectivity in rural areas (25% rural broadband coverage)",
    "💰 High cost of traditional consulting services ($100-500/consultation)",
    "❌ Language barrier - information mostly available in English",
    "🏥 Crop diseases & yields suffer from lack of timely guidance"
])

# Slide 3: Market Opportunity
add_content_slide(prs, "Market Opportunity", [
    "📊 Indian agriculture market: $300+ billion annually",
    "👥 Target market: 50 million small & marginal farmers",
    "📱 Smartphone penetration in rural India: 30%+ and growing",
    "🌾 Digital agriculture adoption rate: 15% (rapidly increasing)",
    "💡 Government push for 'Digital India' & 'Pradhan Mantri e-Kranti'"
])

# Slide 4: Our Solution
add_content_slide(prs, "Our Solution: AI Farming Assistant", [
    "🎤 Voice-based AI system in regional languages (Hindi, Tamil, Telugu, etc.)",
    "⚙️ Runs offline on Raspberry Pi - no internet required",
    "🤖 AI capabilities: Crop disease detection, yield prediction, advisory",
    "🔋 Low power consumption & cost-effective ($30-50 per unit)",
    "✅ Open-source, scalable & adaptable for different crops"
])

# Slide 5: Revenue Streams
add_content_slide(prs, "Business Model & Revenue Streams", [
    "💰 Freemium Model: Basic advisory free, premium features ($2-5/month)",
    "🤝 B2G Partnerships: Government farming schemes (PMKSY reach: 50M farmers)",
    "🏢 Farm Cooperative Integration: Revenue share with agricultural cooperatives",
    "📊 Data Analytics: Anonymized insights sold to agritech companies",
    "🌐 Licensing: Licensing AI models to other platforms & services"
])

# Slide 6: Path to Profitability
add_content_slide(prs, "Sustainability & Path to Profitability", [
    "Year 1: Deploy 10,000 units → ₹20 lakh revenue (pilot phase)",
    "Year 2: Scale to 100,000 units → ₹2.5 crore (government partnerships)",
    "Year 3: Reach 500,000 units → ₹12.5 crore + data licensing",
    "📈 Unit economics: $10 cost → $30-50 lifetime value (4-5x margin)",
    "✅ Breakeven at 50,000 active users (achievable by Month 18)"
])

# Slide 7: Competitive Advantages
add_content_slide(prs, "Competitive Advantages", [
    "🌍 Only offline-first solution (competitors need internet)",
    "🗣️ Multilingual & voice-first (vs. text-based apps)",
    "🤖 AI runs ON-device (Phi-3 3.8B model on CPU - no GPU needed)",
    "💚 Open-source & zero licensing costs (vs. proprietary solutions)",
    "👨‍🌾 Built specifically for Indian farmer context & needs"
])

# Slide 8: Market Entry Strategy
add_content_slide(prs, "Go-to-Market Strategy", [
    "Phase 1 (Months 1-6): Pilot in 3 districts (Tamil Nadu, AP, Punjab)",
    "Phase 2 (Months 6-12): Partner with KRISHI VIGYAN KEARs (KVKs)",
    "Phase 3 (Year 2): State government partnerships via PMKSY scheme",
    "Phase 4 (Year 2-3): Expand to other states via farm cooperatives",
    "🎯 Target: 1 million farmers by end of Year 3"
])

# Slide 9: Financial Projections
add_content_slide(prs, "Financial Projections (3-Year)", [
    "Year 1: Revenue ₹20L | Expenses ₹40L | Loss ₹20L (R&D phase)",
    "Year 2: Revenue ₹2.5Cr | Expenses ₹1.8Cr | Profit ₹70L",
    "Year 3: Revenue ₹12.5Cr | Expenses ₹7Cr | Profit ₹5.5Cr",
    "Breakeven by Month 18 | ROI targets: 3x by Year 3",
    "Capital required: ₹50L for pilot + 1st year operations"
])

# Slide 10: Team & Partnerships
add_content_slide(prs, "Team & Strategic Partnerships", [
    "👨‍💻 Core Team: AI/ML engineers + Agricultural domain experts",
    "🤝 Academic Partner: VIT Chennai (research & R&D)",
    "🏛️ Government Support: Potential NRLM, PMKSY integration",
    "🌱 Partner Organizations: State agricultural departments & KVKs",
    "💼 Advisors: Agritech entrepreneurs & agricultural scientists"
])

# Slide 11: Risk Mitigation
add_content_slide(prs, "Risk Mitigation Strategy", [
    "🔴 Tech Risk: Use proven open-source models (Whisper, Phi3, eSpeak)",
    "📡 Connectivity: Build for offline-first, no internet dependency",
    "💰 Market Risk: Government partnerships reduce revenue volatility",
    "👥 Adoption Risk: Conduct farmer education & demo programs",
    "⚖️ Regulatory: Comply with data privacy & agricultural regulations"
])

# Slide 12: Impact & Vision
add_content_slide(prs, "Social Impact & Vision", [
    "🌾 Empower 10 million Indian farmers with AI advisory by 2030",
    "📈 Increase average farm productivity by 20-30%",
    "💰 Reduce farmer losses from disease/poor practices by ₹5000+ per unit annually",
    "🌍 Model sustainable agritech that works without internet",
    "🚀 Building India's indigenous AI+Agriculture ecosystem"
])

# Slide 13: Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRAND_COLOR

# Main text
cta_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
cta_frame = cta_box.text_frame
cta_frame.text = "Join Us in Revolutionizing\nIndian Agriculture"
p = cta_frame.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Let's Build the Future of Farming Together"
p = subtitle_frame.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_COLOR
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r'c:\Users\thame\Desktop\SEM-8\IOT Domain Analyst\Project\AI_Farming_Assistant_Startup_Pitch.pptx'
prs.save(output_path)
print("✅ PowerPoint presentation created successfully!")
print(f"📁 Saved as: {output_path}")
